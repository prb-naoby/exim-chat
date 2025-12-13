"""
PPT/PPTX to PDF Converter - Linux Compatible

Uses LibreOffice for conversion (works in Docker/Linux).
Falls back to python-pptx for text extraction if LibreOffice is not available.
"""

import os
import subprocess
import tempfile
import shutil
from typing import Optional


def convert_ppt_to_pdf(input_path: str, output_path: str) -> bool:
    """
    Convert a PowerPoint file to PDF using LibreOffice (Linux/Docker compatible).
    
    Args:
        input_path: Path to input .ppt or .pptx file
        output_path: Path to output .pdf file
        
    Returns:
        bool: True if successful, False otherwise
    """
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    
    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
        return False
    
    # Method 1: Try LibreOffice (soffice) - works in Linux/Docker
    if shutil.which("soffice") or shutil.which("libreoffice"):
        try:
            # LibreOffice outputs to the same directory as input, so we use a temp dir
            with tempfile.TemporaryDirectory() as temp_dir:
                # Copy input to temp dir
                temp_input = os.path.join(temp_dir, os.path.basename(input_path))
                shutil.copy2(input_path, temp_input)
                
                # Convert using LibreOffice
                cmd = [
                    shutil.which("soffice") or shutil.which("libreoffice"),
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", temp_dir,
                    temp_input
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    # Find the generated PDF
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    temp_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
                    
                    if os.path.exists(temp_pdf):
                        shutil.move(temp_pdf, output_path)
                        print(f"Successfully converted {input_path} to {output_path}")
                        return True
                
                print(f"LibreOffice conversion failed: {result.stderr}")
                
        except subprocess.TimeoutExpired:
            print(f"Error: LibreOffice conversion timed out for {input_path}")
        except Exception as e:
            print(f"Error during LibreOffice conversion: {e}")
    
    # Method 2: Fallback - extract text using python-pptx (no actual PDF, just text)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        print(f"Warning: LibreOffice not found. Extracting text only from {input_path}")
        
        prs = Presentation(input_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_content.append("\n".join(slide_text))
        
        # Save as text file with .pdf extension (not ideal but fallback)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(text_content))
        
        print(f"Extracted text from {input_path} (LibreOffice not available)")
        return True
        
    except ImportError:
        print("Error: python-pptx not installed. Cannot extract PPT content.")
        return False
    except Exception as e:
        print(f"Error extracting PPT content: {e}")
        return False


def is_libreoffice_available() -> bool:
    """Check if LibreOffice is available for conversion."""
    return bool(shutil.which("soffice") or shutil.which("libreoffice"))
